import os
from pptx import Presentation
import openai

# Set your OpenAI API key
openai.api_key = 'api-open-key'

def create_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Slide layout with title and content
    slide = prs.slides.add_slide(slide_layout)
    
    # Set slide title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Set slide content
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.text = content

def generate_content(prompt):
    # Generate text content using GPT-3.5-turbo or GPT-4
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",  # You can change to "gpt-4" if you have access
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ],
        max_tokens=200,
        n=1,
        temperature=0.7,
    )
    generated_text = response.choices[0].message['content'].strip()
    return generated_text

def main():
    print("Initializing AI model... This may take a moment.")
    
    # Get presentation topic from user
    presentation_topic = input("Enter your presentation topic: ")
    
    # Generate presentation title
    title_prompt = f"Create a title for a presentation about {presentation_topic}:"
    presentation_title = generate_content(title_prompt).split('\n')[0]
    
    # Create a new presentation
    prs = Presentation()
    
    # Add title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = presentation_title
    
    # Generate and add content slides
    slides = []
    for i in range(3):  # Create 3 content slides
        slide_prompt = f"Create a slide title and content for a presentation about {presentation_topic}. Slide {i+1}:"
        slide_content = generate_content(slide_prompt)
        
        lines = slide_content.split('\n')
        slide_title = lines[0].strip()
        slide_content = '\n'.join(lines[1:]).strip()
        
        slides.append((slide_title, slide_content))
    
    for title, content in slides:
        create_slide(prs, title, content)
    
    # Save the presentation
    prs.save('generated_presentation.pptx')
    print(f"Presentation saved as 'generated_presentation.pptx' in {os.getcwd()}")

if __name__ == "__main__":
    main()

# To Run
# install pip install python-pptx openai

